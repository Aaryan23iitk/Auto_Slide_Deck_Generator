"""
Auto Slide Deck Generator
-------------------------
Author: Aaryan Tiwari (BT ME IITK)

What this script does
- Prompts for a topic (or use --topic)
- Uses DuckDuckGo to fetch fresh web snippets
- Calls an OpenAI LLM to synthesize a 7‑slide outline (Title, Overview, 3–4 Key Points, Conclusion)
- Builds a .pptx deck with python-pptx

Quick start
-----------
1) Python 3.10+
2) pip install -U ddgs python-pptx openai tiktoken tenacity
3) Set your OpenAI API key:
   - Linux/macOS: export OPENAI_API_KEY=sk-...
   - Windows (CMD): set OPENAI_API_KEY=sk-...
   - Windows (PowerShell): $Env:OPENAI_API_KEY = "sk-..."
4) Run:
   python auto_slide_deck.py --topic "Climate Change" --outfile "Climate_Change.pptx"

Notes
-----
- The script asks OpenAI to return strict JSON and validates it before slide creation.
- If JSON is malformed, a light repair pass tries to fix it.
- The deck uses simple, clean, bullet-based slides (no external template).

"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
import textwrap
from dataclasses import dataclass
from typing import Any, Dict, List, Tuple
import random
from pptx.dml.color import RGBColor

from tenacity import retry, stop_after_attempt, wait_exponential

# --- Web search (DuckDuckGo) ---
try:
    from ddgs import DDGS  # pip install ddgs
except Exception as e:
    DDGS = None  

# --- OpenAI ---
try:
    from openai import OpenAI  # pip install openai
except Exception as e: 
    OpenAI = None  

# --- PPTX ---
try:
    from pptx import Presentation  # pip install python-pptx
    from pptx.util import Inches, Pt
except Exception as e:  
    Presentation = None  
    Inches = None 
    Pt = None  


# =============================
# Data structures & utilities
# =============================

@dataclass
class WebResult:
    title: str
    snippet: str
    url: str


def sanitize_filename(name: str) -> str:
    name = name.strip().replace(" ", "_")
    return re.sub(r"[^A-Za-z0-9_.-]", "", name)


def clamp(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3] + "..."


# =============================
# Part 2: Web Search Module
# =============================


def get_web_results(query: str, max_results: int = 8) -> Tuple[List[WebResult], str]:
    """Fetch top results and return (list, concatenated_context).
    The concatenated context is trimmed to ~6000 chars to control token cost.
    """
    if DDGS is None:
        raise RuntimeError(
            "duckduckgo-search is not installed."
        )

    results: List[WebResult] = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=max_results):
            title = r.get("title", "").strip()
            snippet = r.get("body", "").strip()
            url = r.get("href", "").strip() or r.get("url", "").strip()
            if title or snippet:
                results.append(WebResult(title=title, snippet=snippet, url=url))

    # Build a compact context block for the LLM
    chunks: List[str] = []
    for i, r in enumerate(results, start=1):
        chunks.append(
            f"[{i}] {r.title}\nSummary: {r.snippet}\nURL: {r.url}\n"
        )
    context = "\n".join(chunks)
    context = clamp(context, 6000)
    return results, context



# =============================
# Part 3: LLM Content Generation
# =============================

SLIDE_SCHEMA_EXAMPLE = {
    "slides": [
        {
            "title": "<Slide Title>",
            "bullets": ["bullet 1", "bullet 2", "bullet 3"],
            "notes": "Presenter notes or empty string"
        }
    ]
}


LLM_SYSTEM_PROMPT = (
    "You are a precise technical writer who produces clean, presentation-ready content. "
    "Return STRICT JSON with keys: slides -> list[ {title, bullets, notes} ]. "
    "Do not include any commentary outside JSON. Keep bullets concise (<= 17 words each). "
    "Cite facts conservatively and avoid unverifiable claims."
)


def format_user_prompt(topic: str, web_context: str, require_seven_slides: bool = True) -> str:
    target_count = 7 if require_seven_slides else 6
    template = f"""
Create a concise slide deck on the topic: \"{topic}\" using BOTH prior knowledge and the provided web snippets.

Required slides ({target_count} total):
1. Title (only the topic as title, no bullets)
2. Overview (5-7 bullets)
3–6.- Key Points / Trends / Arguments (each 5-7 bullets) 
    Thematic sections with clear, specific titles (NOT "Key Point 1" etc.)
   - Titles should summarize the content (e.g., "Trends ","Arguments")
   
7. Conclusion / Takeaways (5-7 bullets)

Guidelines:
- Use plain language; avoid marketing fluff.
- Bullets should be short, scannable, and factual.
- Prefer recent insights inferred from snippets when relevant.
- If the web data seems conflicting, summarize the consensus.

Here are the web snippets (search results):
---
{web_context}
---

Return JSON ONLY in the form:
{json.dumps(SLIDE_SCHEMA_EXAMPLE, ensure_ascii=False)}
"""
    return textwrap.dedent(template).strip()


@retry(wait=wait_exponential(multiplier=1, min=1, max=8), stop=stop_after_attempt(3))
def call_openai_json(prompt: str, model: str) -> Dict[str, Any]:
    if OpenAI is None:
        raise RuntimeError(
            "openai package not installed."
        )
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY environment variable not set.")

    client = OpenAI()
    resp = client.chat.completions.create(
        model=model,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": LLM_SYSTEM_PROMPT},
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
    )
    content = resp.choices[0].message.content
    if not content:
        raise RuntimeError("Empty response from LLM")
    return json.loads(content)


def validate_slide_json(data: Dict[str, Any]) -> List[Dict[str, Any]]:
    # Ensure data has the structure we expect
    if not isinstance(data, dict) or "slides" not in data:
        raise ValueError("LLM did not return a 'slides' list")
    slides = data["slides"]
    if not isinstance(slides, list) or not slides:
        raise ValueError("'slides' must be a non-empty list")

    norm: List[Dict[str, Any]] = []
    for s in slides:
        title = str(s.get("title", "")).strip()
        bullets = s.get("bullets", [])
        notes = str(s.get("notes", "")).strip()
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []
        # Clean bullets (no empties, trim)
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        norm.append({"title": title, "bullets": bullets, "notes": notes})
    return norm


# =============================
# Part 4: PowerPoint Generation
# =============================

EMOJIS = ["✅","📈","🔑","📊","💡"]

def _add_bullets(text_frame, bullets: List[str], slide_index: int = 0):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    emoji = EMOJIS[slide_index % len(EMOJIS)]

    p = text_frame.paragraphs[0]
    p.text = ""
    for i, b in enumerate(bullets):
        bullet_text = f"{emoji} {b}"

        if i == 0:
            p.text = bullet_text
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(50, 50, 50)
        else:
            run = text_frame.add_paragraph()
            run.text = bullet_text
            run.level = 0
            for r in run.runs:
                r.font.size = Pt(20)
                r.font.color.rgb = RGBColor(50, 50, 50)



COLOR_PALETTE = [
    (230, 240, 250),  # light blue
    (245, 245, 245),  # light gray
    (250, 240, 230),  # soft peach
    (240, 250, 240),  # light green
    (255, 255, 240),  # ivory
]

def set_slide_background(slide, rgb_color=None):
    """Apply a random or chosen background color to a slide"""
    fill = slide.background.fill
    fill.solid()
    if rgb_color is None:
        rgb_color = random.choice(COLOR_PALETTE)
    fill.fore_color.rgb = RGBColor(*rgb_color)

def create_pptx(slides: List[Dict[str, Any]], outfile: str) -> str:
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is not installed."
        )
    prs = Presentation()

    # Title slide layout is usually index 0, Title & Content is index 1
    TITLE_LAYOUT = 0
    CONTENT_LAYOUT = 1

    for idx, s in enumerate(slides):
        title = s.get("title") or f"Slide {idx+1}"
        bullets = s.get("bullets", [])
        notes = s.get("notes", "")

        if idx == 0 and bullets == []:
            # Create the title slide
            slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
            title_shape = slide.shapes.title
            title_shape.text = title

            # --- Background color (dark blue) ---
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(25, 50, 100)

            # --- Style the title (center, big, white) ---
            from pptx.enum.text import PP_ALIGN
            paragraph = title_shape.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102) 

            # --- Add subtitle/tagline below title ---
            left = Inches(5)  
            top = Inches(6.5)  
            width = Inches(3)  
            height = Inches(0.5)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_tf = subtitle_box.text_frame
            subtitle_tf.text = "Auto-Generated Deck • Powered by AI"
            for run in subtitle_tf.paragraphs[0].runs:
                run.font.size = Pt(20)
                run.font.italic = True
                run.font.color.rgb = RGBColor(200, 200, 200)

            # --- Add a gold accent line under the title ---
                from pptx.enum.shapes import MSO_SHAPE

                # get title position & size
                title_left = title_shape.left
                title_top = title_shape.top
                title_width = title_shape.width
                title_height = title_shape.height

                # place the line just below the title (with a small gap)
                line_top = title_top + title_height + Inches(0.1)
                line_height = Inches(0.15)

                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    title_left, line_top,
                    title_width, line_height
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(138, 43, 226)  # Purple
                line.line.fill.background()  # remove border


        else:
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
            title_shape = slide.shapes.title
            title_shape.text = title

            # Styling headings
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)  # navy blue

            # underline the title (skip first slide)
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.underline = True

            body = slide.shapes.placeholders[1].text_frame
            if bullets:
                _add_bullets(body, bullets, idx)  
            else:
                body.text = ""

                       
        # Apply background color to every slide here
        set_slide_background(slide)

        # Optional: add presenter notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    safe_out = sanitize_filename(outfile or "deck.pptx")
    if not safe_out.lower().endswith(".pptx"):
        safe_out += ".pptx"
    prs.save(safe_out)
    return safe_out



# =============================
# Part 5: Orchestration
# =============================

def build_deck(topic: str, *,
               model: str,
               max_results: int = 8,
               skip_web: bool = False,
               dry_run: bool = False,
               outfile: str | None = None) -> str | None:
    """End-to-end: search -> LLM -> PPTX.
    Returns output path (or None if dry_run).
    """
    topic = topic.strip()
    if not topic:
        raise ValueError("Topic must not be empty")

    web_ctx = ""
    if not skip_web:
        _, web_ctx = get_web_results(topic, max_results=max_results)

    prompt = format_user_prompt(topic, web_ctx, require_seven_slides=True)
    raw = call_openai_json(prompt, model=model)
    try:
        slides = validate_slide_json(raw)
    except Exception as e:
        # Attempt a light repair if the structure is slightly off
        repaired = {"slides": []}
        if isinstance(raw, list):
            for item in raw:
                if isinstance(item, dict):
                    repaired["slides"].append(item)
        elif isinstance(raw, dict) and "Slides" in raw:
            repaired["slides"] = raw.get("Slides", [])
        else:
            raise
        slides = validate_slide_json(repaired)

    if dry_run:
        print(json.dumps({"slides": slides}, indent=2, ensure_ascii=False))
        return None

    out = outfile or f"{sanitize_filename(topic)}.pptx"
    path = create_pptx(slides, out)
    return path


def parse_args(argv: List[str]) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Auto-generate slide decks with LLM + Web search")
    p.add_argument("--topic", type=str, default=None, help="Topic for the deck")
    p.add_argument("--max-results", type=int, default=8, help="DuckDuckGo results to fetch")
    p.add_argument("--model", type=str, default=os.getenv("OAI_MODEL", "gpt-4o-mini"),
                   help="OpenAI chat model (also via env OAI_MODEL)")
    p.add_argument("--outfile", type=str, default=None, help="Output .pptx filename")
    p.add_argument("--dry-run", action="store_true", help="Print JSON, do not create PPTX")
    p.add_argument("--no-web", action="store_true", help="Skip web search, rely on LLM only")
    return p.parse_args(argv)


def main(argv: List[str] | None = None) -> int:
    ns = parse_args(argv or sys.argv[1:])
    topic = ns.topic or input("Enter topic: ").strip()

    try:
        out = build_deck(
            topic,
            model=ns.model,
            max_results=ns.max_results,
            skip_web=ns.no_web,
            dry_run=ns.dry_run,
            outfile=ns.outfile,
        )
    except Exception as e:
        print(f"[ERROR] {e}")
        return 2

    if out:
        print(f"✅ Deck generated: {out}")
    else:
        print("(dry-run complete; no PPTX created)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())